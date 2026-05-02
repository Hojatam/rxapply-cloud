"""
Daneshyar text extractor.

Usage
-----
  python extract.py <path>
      Detects file type by extension, extracts plain text, prints to stdout.
      Exit 0 = success. Exit non-zero = failure (with stderr message).

Supported formats
-----------------
  .pdf            via pypdf (or PyPDF2 fallback)
  .docx / .doc    via python-docx (.doc not fully supported — convert first)
  .pptx           via python-pptx
  .html / .htm    stdlib html.parser
  .md / .txt      utf-8 read
  .rtf            stdlib (strips control words best-effort)

If a third-party lib is missing, prints a clear "pip install …" hint to
stderr and exits 3 so the server can surface the message to the founder.

Stdout is always UTF-8. Pages / slides are joined with double-newlines so
Daneshyar's parser can split them into separate facts.
"""
import os
import sys
import io
import re

# Force UTF-8 on Windows so Persian/Arabic source text round-trips cleanly.
if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
if hasattr(sys.stderr, "reconfigure"):
    sys.stderr.reconfigure(encoding="utf-8", errors="replace")


def _missing(lib_name: str, install_hint: str) -> None:
    sys.stderr.write(
        f"missing dependency: {lib_name}. Install with:\n  pip install {install_hint}\n"
    )
    sys.exit(3)


def extract_pdf(path: str) -> str:
    try:
        from pypdf import PdfReader  # modern fork (preferred)
    except ImportError:
        try:
            from PyPDF2 import PdfReader  # legacy fallback
        except ImportError:
            _missing("pypdf", "pypdf")
            return ""  # unreachable
    reader = PdfReader(path)
    pages = []
    for i, page in enumerate(reader.pages):
        try:
            txt = page.extract_text() or ""
        except Exception as e:
            sys.stderr.write(f"warn: page {i+1} extraction failed: {e}\n")
            txt = ""
        if txt.strip():
            pages.append(txt.strip())
    return "\n\n".join(pages)


def extract_docx(path: str) -> str:
    try:
        from docx import Document
    except ImportError:
        _missing("python-docx", "python-docx")
        return ""
    doc = Document(path)
    parts = []
    # Body paragraphs.
    for p in doc.paragraphs:
        t = (p.text or "").strip()
        if t:
            parts.append(t)
    # Tables — join rows with " | " so structure is roughly preserved.
    for table in doc.tables:
        for row in table.rows:
            cells = [(c.text or "").strip() for c in row.cells]
            line = " | ".join([c for c in cells if c])
            if line:
                parts.append(line)
    return "\n".join(parts)


def extract_pptx(path: str) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        _missing("python-pptx", "python-pptx")
        return ""
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides):
        bits = [f"Slide {i+1}:"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                t = shape.text.strip()
                if t:
                    bits.append(t)
            # Tables embedded in slides.
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [(c.text or "").strip() for c in row.cells]
                    line = " | ".join([c for c in cells if c])
                    if line:
                        bits.append(line)
        if len(bits) > 1:
            slides.append("\n".join(bits))
    return "\n\n".join(slides)


def extract_html(path: str) -> str:
    # Stdlib only — no BeautifulSoup dependency.
    from html.parser import HTMLParser
    from html import unescape

    class Stripper(HTMLParser):
        def __init__(self):
            super().__init__()
            self.parts = []
            self._skip = 0  # depth inside <script>/<style>
            self._block = False

        def handle_starttag(self, tag, attrs):
            if tag in ("script", "style"):
                self._skip += 1
            elif tag in ("p", "br", "div", "li", "h1", "h2", "h3", "h4",
                          "h5", "h6", "tr", "td", "th"):
                self.parts.append("\n")
                self._block = True

        def handle_endtag(self, tag):
            if tag in ("script", "style") and self._skip > 0:
                self._skip -= 1
            if tag in ("p", "div", "li", "h1", "h2", "h3", "h4", "h5", "h6", "tr"):
                self.parts.append("\n")

        def handle_data(self, data):
            if self._skip > 0:
                return
            t = data.strip()
            if t:
                self.parts.append(t)

    with open(path, "r", encoding="utf-8", errors="replace") as f:
        raw = f.read()
    s = Stripper()
    s.feed(raw)
    text = unescape("".join(s.parts))
    # Collapse runs of blank lines.
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def extract_text_file(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


def extract_rtf(path: str) -> str:
    # Best-effort RTF stripping without external deps.
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        rtf = f.read()
    # Remove RTF groups, control words, hex escapes.
    s = re.sub(r"\\par[d]?", "\n", rtf)
    s = re.sub(r"\\'[0-9a-fA-F]{2}", "", s)
    s = re.sub(r"\\[a-zA-Z]+-?\d* ?", "", s)
    s = re.sub(r"[{}]", "", s)
    return s.strip()


EXTRACTORS = {
    ".pdf":  extract_pdf,
    ".docx": extract_docx,
    ".doc":  extract_docx,   # often works; fails clearly if not
    ".pptx": extract_pptx,
    ".html": extract_html,
    ".htm":  extract_html,
    ".md":   extract_text_file,
    ".txt":  extract_text_file,
    ".rtf":  extract_rtf,
}


def main():
    if len(sys.argv) < 2:
        sys.stderr.write("usage: extract.py <path>\n")
        sys.exit(2)
    path = sys.argv[1]
    if not os.path.isfile(path):
        sys.stderr.write(f"not a file: {path}\n")
        sys.exit(2)
    ext = os.path.splitext(path)[1].lower()
    fn = EXTRACTORS.get(ext)
    if not fn:
        sys.stderr.write(f"unsupported extension: {ext}\n")
        sys.exit(2)
    try:
        text = fn(path)
    except SystemExit:
        raise
    except Exception as e:
        sys.stderr.write(f"extract failed: {type(e).__name__}: {e}\n")
        sys.exit(4)
    if not text or not text.strip():
        sys.stderr.write("warning: no text extracted (file may be image-only or empty)\n")
        # Still exit 0 so caller can decide what to do; print empty stdout.
    sys.stdout.write(text)


if __name__ == "__main__":
    main()
